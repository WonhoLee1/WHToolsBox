"""WHTools Drop Simulator — 소개 프레젠테이션 생성 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
from pathlib import Path

# ── 색상 팔레트 ──────────────────────────────────────────────
C_BG       = RGBColor(0x1a, 0x1a, 0x2e)   # 딥 네이비
C_ACCENT   = RGBColor(0x16, 0x21, 0x3e)   # 패널 배경
C_BLUE     = RGBColor(0x0f, 0x8b, 0x8d)   # 청록 강조
C_GOLD     = RGBColor(0xe9, 0xc4, 0x6a)   # 골드
C_WHITE    = RGBColor(0xff, 0xff, 0xff)
C_LIGHT    = RGBColor(0xd0, 0xe8, 0xff)
C_GRAY     = RGBColor(0x88, 0x99, 0xaa)
C_GREEN    = RGBColor(0x2a, 0x9d, 0x8f)
C_RED      = RGBColor(0xe7, 0x63, 0x5a)

LOGO = Path(r"C:\Users\GOODMAN\WHToolsBox\TVPackageMotionSim\resources\sidebar_logo.png")
BANNER = Path(r"C:\Users\GOODMAN\WHToolsBox\TVPackageMotionSim\ui_banner.png")
W, H = Inches(13.33), Inches(7.5)   # 와이드 16:9

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def fill_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape

def add_text(slide, text, l, t, w, h, size=24, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_bullet_box(slide, items, l, t, w, h, title=None,
                   bg=C_ACCENT, bullet="▸", size=16, title_size=18):
    add_rect(slide, l, t, w, h, bg)
    y = t + Inches(0.18)
    if title:
        add_text(slide, title, l + Inches(0.2), y, w - Inches(0.3),
                 Inches(0.4), size=title_size, bold=True, color=C_GOLD)
        y += Inches(0.42)
    for item in items:
        add_text(slide, f"{bullet}  {item}", l + Inches(0.2), y,
                 w - Inches(0.35), Inches(0.36), size=size, color=C_LIGHT)
        y += Inches(0.34)

def section_bar(slide, text, color=C_BLUE):
    add_rect(slide, 0, Inches(6.9), W, Inches(0.6), color)
    add_text(slide, text, Inches(0.3), Inches(6.92), Inches(8), Inches(0.5),
             size=13, color=C_WHITE)

# ════════════════════════════════════════════════════════════
# Slide 1 — Title
# ════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_BG)

# 좌측 그라데이션 패널
add_rect(sl, 0, 0, Inches(5.5), H, C_ACCENT)
add_rect(sl, 0, 0, Inches(0.12), H, C_BLUE)

# 로고
if LOGO.exists():
    sl.shapes.add_picture(str(LOGO), Inches(1.0), Inches(0.6), Inches(3.5), Inches(3.5))

add_text(sl, "WHTools", Inches(0.4), Inches(4.3), Inches(4.8), Inches(0.8),
         size=44, bold=True, color=C_GOLD)
add_text(sl, "Drop Simulator", Inches(0.4), Inches(5.0), Inches(4.8), Inches(0.7),
         size=32, bold=True, color=C_WHITE)
add_text(sl, "TV Package Motion Simulation Suite", Inches(0.4), Inches(5.7),
         Inches(4.8), Inches(0.45), size=16, color=C_GRAY)
add_text(sl, "Powered by  MuJoCo · OpenRadioss · PySide6",
         Inches(0.4), Inches(6.25), Inches(4.8), Inches(0.4), size=13, color=C_BLUE)

# 우측 배너
if BANNER.exists():
    sl.shapes.add_picture(str(BANNER), Inches(5.8), Inches(1.5), Inches(7.0), Inches(3.1))

add_text(sl, "낙하 시뮬레이션 자동화 플랫폼", Inches(5.8), Inches(4.9),
         Inches(7.0), Inches(0.5), size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl,
    "포장재 설계 단계에서 낙하 충격을 정밀 예측하여\n개발 기간 단축 및 시험 비용 절감",
    Inches(5.8), Inches(5.4), Inches(7.0), Inches(0.8), size=14,
    color=C_GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# Slide 2 — 개요 / 목적
# ════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_BG)
add_rect(sl, 0, 0, W, Inches(1.1), C_ACCENT)
add_rect(sl, 0, 0, Inches(0.12), H, C_BLUE)
add_text(sl, "개요 및 목적", Inches(0.4), Inches(0.2), Inches(10), Inches(0.7),
         size=32, bold=True, color=C_GOLD)

# 3-column cards
card_w = Inches(3.9)
card_h = Inches(4.8)
card_y = Inches(1.3)
cards = [
    ("🎯  배경 / 문제",
     ["TV 세트 낙하 시험은 물리적 시제품 필요",
      "시제품 제작 → 비용·기간 과다",
      "설계 변경 시 반복 시험 부담",
      "포장재 최적화에 정량적 데이터 부족"]),
    ("💡  솔루션",
     ["MuJoCo 기반 실시간 물리 시뮬레이션",
      "OpenRadioss FEM 해석으로 응력 정밀도 확보",
      "낙하 각도·높이 자동 배치 (22가지 시나리오)",
      "배치 실행 → 결과 자동 수집"]),
    ("📈  기대 효과",
     ["시제품 없이 설계 단계 평가 가능",
      "포장재 두께·재질 최적화 가속",
      "시험 횟수 대폭 감소 (비용 절감)",
      "시나리오별 결과 DB 자동 구축"]),
]
for i, (title, bullets) in enumerate(cards):
    x = Inches(0.4) + i * (card_w + Inches(0.15))
    add_bullet_box(sl, bullets, x, card_y, card_w, card_h,
                   title=title, title_size=17, size=15)

section_bar(sl, "WHTools Drop Simulator  ·  개요")

# ════════════════════════════════════════════════════════════
# Slide 3 — 시스템 아키텍처
# ════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_BG)
add_rect(sl, 0, 0, W, Inches(1.1), C_ACCENT)
add_rect(sl, 0, 0, Inches(0.12), H, C_BLUE)
add_text(sl, "시스템 아키텍처", Inches(0.4), Inches(0.2), Inches(10), Inches(0.7),
         size=32, bold=True, color=C_GOLD)

# Layer boxes
layers = [
    ("UI Layer  (PySide6)",   "Control Panel · Model Setup · Batch Launcher · Result Viewer",  C_BLUE),
    ("Simulation Layer",      "WHTs Engine (MuJoCo)  ·  JAX SSR  ·  Physics Solver",           C_GREEN),
    ("FEM Layer",             "OpenRadioss Builder  ·  INP/RAD 변환  ·  Penetration Check",    RGBColor(0x8a,0x63,0xd2)),
    ("Post-Processing Layer", "Multi-Postprocessor  ·  Analysis Pipeline  ·  Reporting",       RGBColor(0xd2,0x7c,0x33)),
]
for i, (name, desc, color) in enumerate(layers):
    y = Inches(1.3) + i * Inches(1.28)
    add_rect(sl, Inches(0.5), y, W - Inches(1.0), Inches(1.1), color)
    add_text(sl, name, Inches(0.7), y + Inches(0.1), Inches(4), Inches(0.45),
             size=18, bold=True, color=C_WHITE)
    add_text(sl, desc, Inches(0.7), y + Inches(0.52), Inches(12.0), Inches(0.45),
             size=14, color=C_LIGHT)

section_bar(sl, "WHTools Drop Simulator  ·  아키텍처")

# ════════════════════════════════════════════════════════════
# Slide 4 — MuJoCo 시뮬레이션
# ════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_BG)
add_rect(sl, 0, 0, W, Inches(1.1), C_ACCENT)
add_rect(sl, 0, 0, Inches(0.12), H, C_BLUE)
add_text(sl, "물리 시뮬레이션 엔진  —  MuJoCo", Inches(0.4), Inches(0.2),
         Inches(12), Inches(0.7), size=30, bold=True, color=C_GOLD)

add_bullet_box(sl,
    ["TV 세트 + 포장재를 MuJoCo MJCF XML로 모델링",
     "낙하 높이 / 각도를 초기 속도·자세로 자동 변환",
     "실시간 뷰어(OpenGL) 또는 오프스크린 렌더러 선택",
     "시뮬레이션 결과를 .pkl 파일로 저장 → 후처리 연계"],
    Inches(0.4), Inches(1.25), Inches(6.0), Inches(2.5),
    title="▶  동작 원리", title_size=17, size=15)

add_bullet_box(sl,
    ["Discrete Block Model : 3×3×3, 3×3×1(Weld) 등 구조 프리셋",
     "Normal / Fast / Rough 3가지 블록 물성 프리셋",
     "Ground Friction 별도 다이얼로그로 정밀 설정",
     "22가지 낙하 시나리오 자동 배치 (IEC/ISTA 기준)"],
    Inches(0.4), Inches(3.9), Inches(6.0), Inches(2.6),
    title="▶  모델 설정", title_size=17, size=15)

add_bullet_box(sl,
    ["멀티 시나리오 병렬 실행 (worker 수 지정)",
     "각 시나리오별 MP4 동영상 자동 캡처 (imageio)",
     "진행률 실시간 모니터링 (Rich 콘솔)",
     "배치 완료 후 결과 폴더 자동 오픈"],
    Inches(6.7), Inches(1.25), Inches(6.3), Inches(4.3),
    title="▶  배치 실행", title_size=17, size=15)

section_bar(sl, "WHTools Drop Simulator  ·  MuJoCo 시뮬레이션")

# ════════════════════════════════════════════════════════════
# Slide 5 — OpenRadioss FEM 해석
# ════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_BG)
add_rect(sl, 0, 0, W, Inches(1.1), C_ACCENT)
add_rect(sl, 0, 0, Inches(0.12), H, C_BLUE)
add_text(sl, "FEM 해석 연동  —  OpenRadioss", Inches(0.4), Inches(0.2),
         Inches(12), Inches(0.7), size=30, bold=True, color=C_GOLD)

add_bullet_box(sl,
    ["MuJoCo 자세(위치·회전)를 Radioss 초기 조건으로 자동 변환",
     "두 가지 변환 모드 : 'parts' (세트 이동) / 'ground' (지면 이동)",
     "Ground Penetration 사전 검사 (0.5 mm 클리어런스)",
     "침투 감지 시 z축 방향으로 자동 보정 후 해석 시작",
     "INP → RAD 자동 변환, 서브루틴 자동 패치"],
    Inches(0.4), Inches(1.25), Inches(6.2), Inches(3.2),
    title="▶  모델 빌더", title_size=17, size=15)

add_bullet_box(sl,
    ["OpenRadioss Starter + Engine 순차 실행",
     "stdout/stderr 실시간 스트리밍 (\\r 정규화 처리)",
     "콘솔 폭 동적 감지 (shutil.get_terminal_size)",
     "해석 완료 후 d3plot / binout 결과 자동 수집"],
    Inches(0.4), Inches(4.55), Inches(6.2), Inches(2.0),
    title="▶  실행 & 모니터링", title_size=17, size=15)

add_bullet_box(sl,
    ["응력(Von Mises) 분포 시각화",
     "에너지 / 변위 이력 그래프",
     "포장재 부품별 최대 응력 비교 테이블",
     "시나리오간 결과 비교 (Multi-Postprocessor)",
     "Excel / CSV 리포트 자동 생성"],
    Inches(6.7), Inches(1.25), Inches(6.3), Inches(5.3),
    title="▶  후처리 & 리포팅", title_size=17, size=15)

section_bar(sl, "WHTools Drop Simulator  ·  OpenRadioss FEM")

# ════════════════════════════════════════════════════════════
# Slide 6 — UI / UX
# ════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_BG)
add_rect(sl, 0, 0, W, Inches(1.1), C_ACCENT)
add_rect(sl, 0, 0, Inches(0.12), H, C_BLUE)
add_text(sl, "사용자 인터페이스", Inches(0.4), Inches(0.2), Inches(12), Inches(0.7),
         size=30, bold=True, color=C_GOLD)

panels = [
    ("🖥  Control Panel",
     ["낙하 시나리오 체크박스 (Select All / Deselect All)",
      "모델 설정 다이얼로그 (블록·마찰·회전 등)",
      "배치 실행 진행률 + 로그 창 (Consolas 9pt)",
      "View 메뉴 → About (빌드 번호 표시)"]),
    ("⚙️  Model Setup Dialog",
     ["Config 트리 : 카테고리별 파라미터 관리",
      "📂/📁 Expand·Fold (선택 항목 / 전체)",
      "Normal / Fast / Rough 블록 프리셋 버튼",
      "Ground Friction 전용 설정 다이얼로그"]),
    ("📊  Result Viewer",
     ["시나리오별 응력·에너지 그래프 오버레이",
      "부품별 컬러 맵 렌더링",
      "MP4 동영상 재생 & 저장",
      "Excel 리포트 원클릭 출력"]),
]
for i, (title, bullets) in enumerate(panels):
    x = Inches(0.4) + i * Inches(4.25)
    add_bullet_box(sl, bullets, x, Inches(1.25), Inches(4.1), Inches(5.3),
                   title=title, title_size=16, size=14)

section_bar(sl, "WHTools Drop Simulator  ·  UI / UX")

# ════════════════════════════════════════════════════════════
# Slide 7 — 기술 스택
# ════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_BG)
add_rect(sl, 0, 0, W, Inches(1.1), C_ACCENT)
add_rect(sl, 0, 0, Inches(0.12), H, C_BLUE)
add_text(sl, "기술 스택", Inches(0.4), Inches(0.2), Inches(12), Inches(0.7),
         size=30, bold=True, color=C_GOLD)

stacks = [
    ("Physics",    ["MuJoCo 3.x", "JAX / JAXlib", "SciPy (spatial, optimize)"]),
    ("FEM",        ["OpenRadioss", "INP↔RAD 변환기", "NumPy (행렬 연산)"]),
    ("UI",         ["PySide6 (Qt6)", "PyQtGraph", "PyVista (3D)"]),
    ("Video",      ["imageio + imageio-ffmpeg", "mujoco.Renderer (오프스크린)", "libx264 MP4 인코딩"]),
    ("Build",      ["PyInstaller (단일 배포)", "Miniconda vdmc 환경", "Git (빌드 번호 자동 증분)"]),
    ("Logging",    ["Rich (RichHandler)", "shutil.get_terminal_size", "콘솔 폭 동적 적용"]),
]
col_w = Inches(4.1)
row_h = Inches(2.5)
for i, (cat, items) in enumerate(stacks):
    col = i % 3
    row = i // 3
    x = Inches(0.4) + col * (col_w + Inches(0.15))
    y = Inches(1.25) + row * (row_h + Inches(0.12))
    add_bullet_box(sl, items, x, y, col_w, row_h,
                   title=cat, title_size=16, size=14, bullet="◆")

section_bar(sl, "WHTools Drop Simulator  ·  기술 스택")

# ════════════════════════════════════════════════════════════
# Slide 8 — 워크플로우
# ════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_BG)
add_rect(sl, 0, 0, W, Inches(1.1), C_ACCENT)
add_rect(sl, 0, 0, Inches(0.12), H, C_BLUE)
add_text(sl, "업무 워크플로우", Inches(0.4), Inches(0.2), Inches(12), Inches(0.7),
         size=30, bold=True, color=C_GOLD)

steps = [
    ("1\n모델\n로드",    "JSON / PKL\n프로파일 불러오기"),
    ("2\n시나리오\n선택", "낙하 높이·각도\n22가지 선택"),
    ("3\nMuJoCo\n시뮬"),  ("4\nRadioss\n빌드"),
    ("5\nFEM\n해석"),     ("6\n결과\n분석"),
]
steps = [
    ("①  모델 로드",     "TV 세트 + 포장재\nJSON 프로파일"),
    ("②  시나리오 선택", "낙하 높이 / 각도\n22가지 선택"),
    ("③  MuJoCo 실행",  "물리 시뮬레이션\n자세 & 속도 계산"),
    ("④  Radioss 빌드", "FEM 모델 생성\n침투 보정 포함"),
    ("⑤  FEM 해석",     "OpenRadioss\n응력 / 에너지"),
    ("⑥  결과 분석",    "그래프·리포트\nMP4 저장"),
]
box_w = Inches(1.9)
box_h = Inches(2.8)
gap   = Inches(0.22)
start_x = Inches(0.4)
y0 = Inches(1.4)
colors_step = [C_BLUE, C_GREEN, RGBColor(0x8a,0x63,0xd2),
               RGBColor(0xd2,0x7c,0x33), C_RED, C_GOLD]

for i, (title, desc) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    add_rect(sl, x, y0, box_w, box_h, colors_step[i])
    add_text(sl, title, x + Inches(0.1), y0 + Inches(0.15), box_w - Inches(0.15),
             Inches(0.5), size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, desc, x + Inches(0.1), y0 + Inches(0.75), box_w - Inches(0.15),
             Inches(1.8), size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        ax = x + box_w + Inches(0.04)
        add_text(sl, "▶", ax, y0 + Inches(1.1), Inches(0.2), Inches(0.4),
                 size=18, color=C_GRAY, align=PP_ALIGN.CENTER)

add_text(sl,
    "배치 모드에서는 ②~⑤ 단계가 선택된 모든 시나리오에 대해 자동 반복 실행됩니다.",
    Inches(0.4), Inches(4.6), W - Inches(0.8), Inches(0.45),
    size=14, color=C_GRAY, align=PP_ALIGN.CENTER)

section_bar(sl, "WHTools Drop Simulator  ·  워크플로우")

# ════════════════════════════════════════════════════════════
# Slide 9 — 주요 성과 / 차별점
# ════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_BG)
add_rect(sl, 0, 0, W, Inches(1.1), C_ACCENT)
add_rect(sl, 0, 0, Inches(0.12), H, C_BLUE)
add_text(sl, "주요 성과 & 차별점", Inches(0.4), Inches(0.2), Inches(12), Inches(0.7),
         size=30, bold=True, color=C_GOLD)

feats = [
    ("🔗  MuJoCo ↔ Radioss 완전 자동화",
     "자세 변환 · 침투 보정 · 모델 빌드까지 원클릭"),
    ("🎬  시뮬레이션 동영상 자동 캡처",
     "배치 실행 중 시나리오별 MP4 저장 (imageio + libx264)"),
    ("📐  Ground Penetration 사전 보정",
     "0.5 mm 클리어런스 자동 검사 → Radioss 오류 방지"),
    ("🖥  동적 콘솔 폭 적응",
     "shutil.get_terminal_size → 어떤 터미널에서도 깔끔한 출력"),
    ("📦  단일 EXE 배포",
     "PyInstaller : MuJoCo · JAX · Radioss 빌더 포함 All-in-One"),
    ("📊  통합 후처리 & 리포팅",
     "시나리오간 비교 · Excel 리포트 · 3D 응력 맵 자동 생성"),
]
for i, (title, desc) in enumerate(feats):
    col = i % 2
    row = i // 2
    x = Inches(0.4) + col * Inches(6.4)
    y = Inches(1.3) + row * Inches(1.7)
    add_rect(sl, x, y, Inches(6.1), Inches(1.5), C_ACCENT)
    add_rect(sl, x, y, Inches(0.08), Inches(1.5), C_BLUE)
    add_text(sl, title, x + Inches(0.2), y + Inches(0.1), Inches(5.7), Inches(0.5),
             size=16, bold=True, color=C_GOLD)
    add_text(sl, desc, x + Inches(0.2), y + Inches(0.6), Inches(5.7), Inches(0.7),
             size=14, color=C_LIGHT)

section_bar(sl, "WHTools Drop Simulator  ·  성과 & 차별점")

# ════════════════════════════════════════════════════════════
# Slide 10 — 마무리
# ════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, C_BG)
add_rect(sl, 0, 0, W, H, C_ACCENT)
add_rect(sl, 0, 0, Inches(0.15), H, C_GOLD)

if LOGO.exists():
    sl.shapes.add_picture(str(LOGO), Inches(4.4), Inches(0.4), Inches(4.5), Inches(4.5))

add_text(sl, "WHTools Drop Simulator", Inches(0.5), Inches(1.8), Inches(7), Inches(0.9),
         size=36, bold=True, color=C_GOLD)
add_text(sl, "포장 설계 혁신을 위한\n낙하 시뮬레이션 자동화 플랫폼",
         Inches(0.5), Inches(2.75), Inches(7), Inches(1.1), size=22, color=C_WHITE)
add_text(sl, "whbest.lee@gmail.com",
         Inches(0.5), Inches(4.2), Inches(7), Inches(0.45), size=15, color=C_GRAY)

add_text(sl, "Build #115  ·  MuJoCo · OpenRadioss · PySide6",
         Inches(0.5), Inches(6.5), Inches(12), Inches(0.45),
         size=12, color=C_GRAY, align=PP_ALIGN.CENTER)

# ── 저장 ──────────────────────────────────────────────────
out = Path(r"C:\Users\GOODMAN\WHToolsBox\WHTools_DropSimulator_Overview.pptx")
prs.save(str(out))
print(f"Saved → {out}")
