# -*- coding: utf-8 -*-
"""WHTools TV Drop Motion Simulator v6.0 — 프로젝트 상세설명 및 사용 설명서 PPTX 생성 스크립트"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import os

# ── 라이트 테마 색상 팔레트 ──────────────────────────────────────────────
C_BG        = RGBColor(0xf4, 0xf7, 0xfc)   # 소프트 블루그레이 (배경)
C_PRIMARY   = RGBColor(0x1a, 0x36, 0x5d)   # 딥 로열 블루 (타이틀 및 헤더)
C_SECONDARY = RGBColor(0x2b, 0x6c, 0xb0)   # 스카이 블루 (강조선 및 강조 텍스트)
C_CARD_BG   = RGBColor(0xff, 0xff, 0xff)   # 카드 내부 (흰색)
C_ACCENT    = RGBColor(0xeb, 0xf8, 0xff)   # 소프트 블루 (서브 카드 배경)
C_TEXT      = RGBColor(0x2d, 0x37, 0x48)   # 다크 차콜 그레이 (본문 텍스트)
C_LIGHT     = RGBColor(0x4a, 0x55, 0x68)   # 다크 슬레이트 그레이 (부제목 및 서브 텍스트)
C_BORDER    = RGBColor(0xd2, 0xe4, 0xf6)   # 카드 테두리 블루그레이
C_GOLD      = RGBColor(0xdd, 0x6b, 0x20)   # 골드/오렌지 (주의 또는 포인트)
C_GREEN     = RGBColor(0x2f, 0x85, 0x5a)   # 녹색 (성공/이점)
C_RED       = RGBColor(0xc5, 0x30, 0x30)   # 적색 (경고/이슈)
C_GRAY      = RGBColor(0xa0, 0xae, 0xc0)   # 연한 그레이 (구분선 등)
C_WHITE     = RGBColor(0xff, 0xff, 0xff)

LOGO = Path(r"C:\Users\GOODMAN\WHToolsBox\TVPackageMotionSim\resources\sidebar_logo.png")
BANNER = Path(r"C:\Users\GOODMAN\WHToolsBox\TVPackageMotionSim\ui_banner.png")
W, H = Inches(13.33), Inches(7.5)   # 와이드 16:9

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def blank_slide(prs):
    """빈 슬라이드 레이아웃 생성"""
    layout = prs.slide_layouts[6]   # 6: 완전히 빈 레이아웃
    return prs.slides.add_slide(layout)

def fill_bg(slide, color: RGBColor):
    """슬라이드 배경색 채우기"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color: RGBColor, border_color=None):
    """지정된 크기의 사각형 추가"""
    shape = slide.shapes.add_shape(1, l, t, w, h)  # 1 = RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=15, bold=False, color=C_TEXT,
             align=PP_ALIGN.LEFT, wrap=True):
    """지정된 서식으로 텍스트 상자 추가 (Segoe UI 적용 및 다중 행 분할 처리)"""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    
    # 텍스트에 \n이 포함된 경우 각각의 단락(Paragraph)으로 추가하여 텍스트가 겹치지 않게 처리
    lines = text.split('\n')
    is_first = True
    for line in lines:
        if is_first:
            p = tf.paragraphs[0]
            is_first = False
        else:
            p = tf.add_paragraph()
            
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = "Segoe UI"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txb

def add_card_box(slide, items, l, t, w, h, title=None,
                 bg=C_CARD_BG, border=C_BORDER, bullet="•", size=13, title_size=15,
                 title_color=C_PRIMARY, text_color=C_TEXT):
    """내부에 불릿 항목이 들어간 카드 박스 추가 (Segoe UI 및 단일 텍스트 프레임 사용하여 자동 줄바꿈 겹침 방지)"""
    # 배경 사각형 추가
    add_rect(slide, l, t, w, h, bg, border_color=border)
    
    # 하나의 텍스트 박스를 생성하여 PowerPoint 자체 레이아웃 엔진이 자동 줄바꿈 및 줄간격을 계산하게 함
    txb = slide.shapes.add_textbox(l + Inches(0.15), t + Inches(0.15), w - Inches(0.3), h - Inches(0.3))
    tf  = txb.text_frame
    tf.word_wrap = True
    
    # 텍스트 박스 내부 여백을 조절하여 패딩 확보
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    
    is_first = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.size = Pt(title_size)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.space_after = Pt(12)  # 제목 아래 여백 설정
        is_first = False
        
    for item in items:
        if is_first:
            p = tf.paragraphs[0]
            is_first = False
        else:
            p = tf.add_paragraph()
            
        bullet_str = f"{bullet}  " if bullet else ""
        p.text = f"{bullet_str}{item}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.color.rgb = text_color
        p.space_after = Pt(8)  # 불릿 항목 간 세로 여백 설정
        p.line_spacing = 1.15  # 줄간격 1.15배 설정


def section_bar(slide, text, color=C_PRIMARY):
    """하단 정보 바 영역 추가"""
    add_rect(slide, 0, Inches(6.9), W, Inches(0.6), color)
    add_text(slide, text, Inches(0.4), Inches(6.92), Inches(10), Inches(0.55),
             size=12, color=C_WHITE, bold=True)
    add_text(slide, "WHTools TV Drop Motion Simulator v6.0", Inches(9.5), Inches(6.92), Inches(3.4), Inches(0.55),
             size=10, color=C_ACCENT, align=PP_ALIGN.RIGHT)

def add_header(slide, title_text, step_text=None):
    """상단 헤더 표시 영역 추가"""
    # 배경 패널
    add_rect(slide, 0, 0, W, Inches(1.1), C_PRIMARY)
    # 강조 라인
    add_rect(slide, 0, Inches(1.06), W, Inches(0.04), C_SECONDARY)
    
    # 텍스트
    if step_text:
        add_text(slide, step_text, Inches(0.4), Inches(0.15), Inches(8), Inches(0.35),
                 size=12, bold=True, color=C_SECONDARY)
        add_text(slide, title_text, Inches(0.4), Inches(0.45), Inches(10), Inches(0.55),
                 size=28, bold=True, color=C_WHITE)
    else:
        add_text(slide, title_text, Inches(0.4), Inches(0.25), Inches(10), Inches(0.65),
                 size=30, bold=True, color=C_WHITE)

# ==============================================================================
# Slide 1 — Title
# ==============================================================================
sl1 = blank_slide(prs)
fill_bg(sl1, C_BG)

# 좌측 블루 그라데이션 포인트 영역
add_rect(sl1, 0, 0, Inches(5.0), H, C_PRIMARY)
add_rect(sl1, Inches(5.0), 0, Inches(0.12), H, C_SECONDARY)

# 로고 삽입
if LOGO.exists():
    sl1.shapes.add_picture(str(LOGO), Inches(0.8), Inches(0.5), Inches(3.4), Inches(3.4))

add_text(sl1, "WHTools", Inches(0.4), Inches(4.1), Inches(4.2), Inches(0.7),
         size=42, bold=True, color=C_WHITE)
add_text(sl1, "Drop Simulator", Inches(0.4), Inches(4.7), Inches(4.2), Inches(0.6),
         size=30, bold=True, color=C_SECONDARY)
add_text(sl1, "TV Package Motion Simulation Suite", Inches(0.4), Inches(5.4),
         Inches(4.2), Inches(0.45), size=14, color=C_ACCENT)
add_text(sl1, "Powered by  MuJoCo · JAX · OpenRadioss · PySide6",
         Inches(0.4), Inches(6.0), Inches(4.2), Inches(0.4), size=11, color=C_ACCENT)

# 우측 배너 및 타이틀
if BANNER.exists():
    sl1.shapes.add_picture(str(BANNER), Inches(5.6), Inches(1.2), Inches(7.2), Inches(3.2))

add_text(sl1, "프로젝트 상세설명 및 사용 설명서", Inches(5.6), Inches(4.7),
         Inches(7.2), Inches(0.5), size=24, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
add_text(sl1, "MuJoCo 물리 엔진과 JAX 자동 미분을 통합한\nTV 패키지 낙하 시뮬레이션 및 자율 구조 해석 통합 플랫폼",
         Inches(5.6), Inches(5.3), Inches(7.2), Inches(0.7), size=14,
         color=C_TEXT, align=PP_ALIGN.CENTER)

add_text(sl1, "Build #115  ·  WHTools Engineering Team", Inches(5.6), Inches(6.3),
         Inches(7.2), Inches(0.4), size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)


# ==============================================================================
# Slide 2 — 개요 및 목적
# ==============================================================================
sl2 = blank_slide(prs)
fill_bg(sl2, C_BG)
add_header(sl2, "개요 및 목적", "01. PROJECT OVERVIEW")

card_w = Inches(3.9)
card_h = Inches(4.8)
card_y = Inches(1.35)

cards2 = [
    ("⚠️  기존 물리 시험의 한계",
     ["실제 TV 세트 낙하 충격 시험을 위해서는 매회 비싼 실물 시제품 제작이 불가피함",
      "설계가 바뀔 때마다 시제품 재제작과 낙하 시험의 반복으로 비용/기간 과다 소요",
      "복잡한 포장재 구조의 실시간 응력 및 변위 데이터 획득이 매우 어렵고 제한적임"],
     C_CARD_BG, C_RED),
    ("💡  WHTools 통합 솔루션",
     ["MuJoCo 물리 엔진 기반의 초고속(실시간급) 강체-유연체 물리 시뮬레이터 제공",
      "JAX 기반 Kirchhoff 박판 솔버를 탑재해 마커 궤적만으로 자율 구조 해석 연동",
      "OpenRadioss 비선형 유한요소(FEM) 솔버와의 연계로 정밀 해석 데이터 검증"],
     C_ACCENT, C_PRIMARY),
    ("📈  도입 효과 및 기대 가치",
     ["시제품 제작 이전 설계 단계에서 낙하 거동과 충격 강도를 즉시 가상 평가 가능",
      "부품 두께 및 재질 최적화(Cushion, Box)를 통해 개발 비용 70% 이상 절감",
      "22가지 ISTA 표준 시험 규격 만족 여부를 단시간 내 스크리닝하여 신뢰성 확보"],
     C_CARD_BG, C_GREEN),
]

for i, (title, bullets, bg_color, title_color) in enumerate(cards2):
    x = Inches(0.4) + i * (card_w + Inches(0.18))
    add_card_box(sl2, bullets, x, card_y, card_w, card_h,
                 title=title, title_size=15, size=13, title_color=title_color, bg=bg_color)

section_bar(sl2, "WHTools Drop Simulator  ·  개요 및 비즈니스 목적")


# ==============================================================================
# Slide 3 — 시스템 아키텍처 및 데이터 흐름
# ==============================================================================
sl3 = blank_slide(prs)
fill_bg(sl3, C_BG)
add_header(sl3, "시스템 아키텍처 및 데이터 흐름", "02. SYSTEM ARCHITECTURE")

# 가로형 4단계 레이어 배치
layer_w = Inches(2.9)
layer_h = Inches(4.3)
layer_y = Inches(1.5)
gap = Inches(0.18)

layers = [
    ("UI Layer (PySide6)",
     ["Control Center UI Panel",
      "Model Setup Dialog",
      "Result Viewer (Plot/Video)",
      "XML Live Editor",
      "External Tool Config"], C_PRIMARY),
    ("Simulation Layer",
     ["WHTs Engine (MuJoCo)",
      "관성 텐서 자동 보정",
      "Friction/Damping 튜닝",
      "이산 블록 & 용접 프리셋",
      "Trajectory 스냅샷 수집"], C_SECONDARY),
    ("FEM Layer",
     ["OpenRadioss Builder",
      "INP / RAD 파일 자동 변환",
      "지면 침투(Penetration)\n자동 감지 및 Z축 오프셋 보정",
      "Starter/Engine 배치 실행"], C_SECONDARY),
    ("Post-Processing",
     ["JAX Kirchhoff 솔버",
      "Von-Mises 응력장 자동 계산",
      "VTKHDF 시계열 출력 (ParaView)",
      "GLB 3D 메시 익스포트",
      "배치 엑셀 리포트 자동 생성"], C_PRIMARY),
]

for i, (name, items, color) in enumerate(layers):
    x = Inches(0.4) + i * (layer_w + gap)
    add_card_box(sl3, items, x, layer_y, layer_w, layer_h,
                 title=name, title_size=14, size=11, title_color=color, bg=C_CARD_BG)
    if i < 3:
        # 연결 화살표 표시
        ax = x + layer_w + Inches(0.02)
        add_text(sl3, "▶", ax, layer_y + Inches(1.8), Inches(0.15), Inches(0.4),
                 size=14, color=C_SECONDARY, align=PP_ALIGN.CENTER)

# 데이터 흐름 요약 박스
add_rect(sl3, Inches(0.4), Inches(5.95), W - Inches(0.8), Inches(0.8), C_ACCENT, border_color=C_BORDER)
add_text(sl3, "💡  핵심 데이터 파이프라인: 입력 설정(JSON) → 물리 거동(MuJoCo) → 스냅샷 수집 → 자율 해석(JAX Kirchhoff) → 고정밀 FEM 검증(OpenRadioss) & 3D 가시화(ParaView/VTK)",
         Inches(0.6), Inches(6.05), W - Inches(1.2), Inches(0.6), size=12, color=C_TEXT, bold=True)

section_bar(sl3, "WHTools Drop Simulator  ·  시스템 구조 및 파이프라인")


# ==============================================================================
# Slide 4 — 핵심 기능 1: Simulation Control Center UI
# ==============================================================================
sl4 = blank_slide(prs)
fill_bg(sl4, C_BG)
add_header(sl4, "통합 제어 패널 (Simulation Control Center)", "03. KEY FEATURE  —  GUI CONTROL CENTER")

# 좌/우 분할 레이아웃
add_card_box(sl4,
    ["실시간 재생 제어: Play / Pause / Reset 단축키 제공 및 0.1x~10.0x 배속 재생 지원",
     "Timeline 탐색: 타임라인 스냅샷 슬라이더로 충격 전후 프레임을 초정밀(Frame-by-frame)로 분석 가능",
     "카메라 뷰 단축 피팅: +X, -X, +Y, -Y, +Z, -Z, ISO 시점 전환 버튼을 배치하여 최적 관찰 각도 제공",
     "상태 모니터링: 시뮬레이션 물리 시간, 스텝 정보 및 발생한 경고/오류 메시지 직관적 실시간 표출"],
    Inches(0.4), Inches(1.4), Inches(6.0), Inches(5.2),
    title="▶  실시간 시뮬레이션 및 카메라 제어", title_size=16, bg=C_CARD_BG)

add_card_box(sl4,
    ["Config Tree 인터페이스: 치수(Box, Set), 쿠션 강성, 용접 조건 등의 파라미터를 카테고리 트리로 조작",
     "XML Live Editor 내장: MuJoCo XML 파일을 GUI 내부 텍스트 에디터로 직접 수정하고, Apply 시 실시간 재로드",
     "실시간 거동 그래프: 선택된 포장재 및 내부 SET 마커들의 XYZ 위치, 속도 변화 파형을 PyQtGraph 실시간 모니터링",
     "외부 연동 설정: [View] -> [Edit External Tools INI]를 통해 별도 파일 조작 없이 ParaView, Radioss 실행 경로 변경"],
    Inches(6.7), Inches(1.4), Inches(6.2), Inches(5.2),
    title="▶  모델 파라미터 및 외부 도구 관리", title_size=16, bg=C_CARD_BG)

section_bar(sl4, "WHTools Drop Simulator  ·  통합 UI 및 실시간 파라미터 튜닝")


# ==============================================================================
# Slide 5 — 핵심 기능 2: JAX 기반 자율 구조 해석
# ==============================================================================
sl5 = blank_slide(prs)
fill_bg(sl5, C_BG)
add_header(sl5, "JAX 기반 자율 구조 해석 파이프라인", "04. KEY FEATURE  —  DIFFERENTIABLE PHYSICS & FEA")

add_card_box(sl5,
    ["마커 궤적(3D Trajectory) 기반 무치수 해석 (Minimalist): CAD 세부 도면이나 복잡한 전처리(Mesh, Boundary) 없이, 물체 표면에 장착된 마커의 거동 데이터만으로 변형 상태 역추정",
     "JAX 가속화 Kirchhoff 박판 솔버: 고도화된 선형 시스템 솔버를 XLA 컴파일러 및 JAX `vmap`, `jit`으로 고속화하여, 수천 프레임의 구조 해석을 단 몇 초 만에 완료",
     "비선형 Kirchhoff 평판 이론: 대변형(Large Deformation) 거동에 최적화된 변분법 수식을 적용해 박판(Paperbox, Chassis)의 비선형 좌굴 및 충격 응력장 예측"],
    Inches(0.4), Inches(1.4), Inches(6.0), Inches(5.2),
    title="▶  자율 해석 알고리즘 및 엔진 성능", title_size=16, bg=C_CARD_BG)

add_card_box(sl5,
    ["전 파츠 Von-Mises 응력 해석: 포장 박스(Paperbox), 완충재(Cushion), Chassis, OpenCell 등 주요 구성 파츠별 응력(Stress)과 변위(Displacement) 자동 계산",
     "통계적 자율 정렬 (Statistical Alignment): 기준 좌표계나 회전 행렬 정보가 사전에 없는 원시 마커 데이터라도 주성분 분석(PCA) 등 통계적 기법으로 기하 구조 자동 정렬",
     "분석 결과 경량 영구 저장: 최종 계산된 응력/변위 시계열 데이터를 float32 타입으로 최적화하여 `results/latest_results.pkl`에 바이너리로 저장 및 신속 압축"],
    Inches(6.7), Inches(1.4), Inches(6.2), Inches(5.2),
    title="▶  주요 산출물 및 데이터 처리", title_size=16, bg=C_CARD_BG)

section_bar(sl5, "WHTools Drop Simulator  ·  JAX 가속 자율 해석 엔진")


# ==============================================================================
# Slide 6 — 핵심 기능 3: ISTA-6 Amazon 배치 해석
# ==============================================================================
sl6 = blank_slide(prs)
fill_bg(sl6, C_BG)
add_header(sl6, "ISTA-6 Amazon 배치 낙하 해석", "05. KEY FEATURE  —  BATCH SIMULATION")

add_card_box(sl6,
    ["ISTA-6 Amazon 규격 최적화: 물류 유통 환경의 충격을 평가하는 Amazon.com 6-Parcel(ISTA 6-Amazon-SIOC) 및 LTL 낙하 시험 규격을 완벽하게 충족",
     "22가지 낙하 시나리오 구성: 모서리(Corner) 낙하, 에지(Edge) 낙하, 평면(Face) 낙하 등 다양한 각도 및 높이별 시퀀스를 자동으로 인스턴스화",
     "병렬 멀티 프로세싱 (Multi-Workers): 하드웨어 성능에 맞게 CPU 워커(Workers) 개수를 조작하여, 여러 시나리오를 헤드리스 모드로 동시 병렬 실행"],
    Inches(0.4), Inches(1.4), Inches(6.0), Inches(5.2),
    title="▶  배치 해석 시나리오 및 실행 성능", title_size=16, bg=C_CARD_BG)

add_card_box(sl6,
    ["코너 이력 데이터 추출: Cushion, Chassis, OpenCell 등 주요 부위의 8개 코너에 대한 실시간 XYZ 충격 거동 및 최대 속도/가속도 이력 데이터 수집",
     "최적화 연계 파일 자동 출력: 시뮬레이션 실행 완료 시 위상 최적화(Topology Optimization) 스크립트와 연동 가능한 `topo_arg.txt` 파일 자동 생성",
     "Rich 실시간 진행률 대시보드: 콘솔창을 통해 시나리오별 CPU 점유, 완료 시점, 성공 여부, 동영상 캡처 진행 상태를 한눈에 모니터링"],
    Inches(6.7), Inches(1.4), Inches(6.2), Inches(5.2),
    title="▶  결과 수집 및 리포팅 자동화", title_size=16, bg=C_CARD_BG)

section_bar(sl6, "WHTools Drop Simulator  ·  배치 시뮬레이션 및 데이터 수집")


# ==============================================================================
# Slide 7 — 사용 설명서 1: 프로그램 실행 및 기본 설정
# ==============================================================================
sl7 = blank_slide(prs)
fill_bg(sl7, C_BG)
add_header(sl7, "사용 방법  —  프로그램 실행 및 기본 설정", "06. OPERATION MANUAL  —  GETTING STARTED")

add_card_box(sl7,
    ["1. Miniconda 가상환경 활성화 (Python 3.10+ 권장)\n   > conda activate vdmc\n   또는 필요한 의존 패키지(jax, mujoco, pyvista, pyside6) 설치",
     "2. run_drop_simulation_cases_v6.py 내부 설정 편집\n   - 파일 하단의 `if __name__ == '__main__':` 블록 참조\n   - 호출할 셋업 함수 지정 (예: `test_case_1_setup`)\n   - `use_viewer=True` 설정 시 Control Center GUI 활성화",
     "3. 파이썬 진입점 실행\n   > python TVPackageMotionSim/run_drop_simulation_cases_v6.py\n   - 실행 시 MuJoCo Viewer와 Control Center Panel이 동시 구동"],
    Inches(0.4), Inches(1.4), Inches(6.0), Inches(5.2),
    title="▶  프로그램 실행 절차 (Execution Process)", title_size=16, bg=C_CARD_BG, bullet=None)

add_card_box(sl7,
    ["외부 도구 설정 파일: external_tools_config.ini\n- 본 플랫폼은 외부 솔버 및 뷰어와의 고속 연동을 위해 환경 변수 대신 INI 설정 파일을 사용합니다.",
     "주요 설정 파라미터:\n- [Tools] 카테고리 내 실행 파일 절대 경로 작성\n- openradioss_path : OpenRadioss Starter/Engine 실행 파일 경로\n- paraview_path : ParaView.exe 실행 경로\n- lsprepost_path : LS-PrePost.exe 실행 경로",
     "GUI를 통한 원클릭 편집:\n- 상단 메뉴 [View] -> [Edit External Tools Config (INI)] 선택 시 텍스트 에디터 창이 열려 즉시 수정 및 저장 가능"],
    Inches(6.7), Inches(1.4), Inches(6.2), Inches(5.2),
    title="▶  외부 도구 경로 및 연동 설정 (Config INI)", title_size=16, bg=C_ACCENT, bullet="•")

section_bar(sl7, "WHTools Drop Simulator  ·  사용 준비 및 환경 구성")


# ==============================================================================
# Slide 8 — 사용 설명서 2: 모델 설정 및 물리 파라미터 튜닝
# ==============================================================================
sl8 = blank_slide(prs)
fill_bg(sl8, C_BG)
add_header(sl8, "사용 방법  —  모델 설정 및 물리 파라미터 튜닝", "07. OPERATION MANUAL  —  MODEL TUNING")

add_card_box(sl8,
    ["치수 및 강성 파라미터 설정:\n- Model Setup Dialog 내의 Config Tree 메뉴를 확장하여 Box 외곽 치수, 내부 TV SET 치수, Cushion 두께 등을 직관적으로 수정",
     "접촉 및 마찰 계수 조정:\n- 포장 박스 내부 및 낙하 지면(Ground)의 마찰 계수를 전용 Friction 다이얼로그로 정밀 조정하여 충격 반발 형태 튜닝",
     "물성 프리셋 선택:\n- Normal(해석 정밀도와 속도의 균형), Fast(빠른 스크리닝 목적), Rough(초기 레이아웃 확인) 등 3가지 사양 프리셋 버튼 제공"],
    Inches(0.4), Inches(1.4), Inches(6.0), Inches(5.2),
    title="▶  Model Setup Dialog를 통한 간편 튜닝", title_size=16, bg=C_CARD_BG, bullet="▸")

add_card_box(sl8,
    ["관성 텐서 자동 보정 (Assembly Inertia Correction):\n- TV SET 실물의 목표 질량, 무게중심(CoG), 관성모멘트(MoI)를 GUI에 입력 시, MuJoCo XML 파일 내에 가상의 fullinertia 보정 바디를 자동 삽입하여 거동 싱크 정합",
     "이산 블록 및 Weld 연결 (Discrete & Weld Blocks):\n- Cushion 등 충격을 많이 받는 파츠를 여러 개의 이산화된 블록(예: 3x3x3 등) 구조로 자동 분할 생성\n- 각 블록 접합부에 solref/solimp 기반 용접(Weld) 강도를 적용하여 한계 하중 도달 시 파손 및 탈거 거동 구현"],
    Inches(6.7), Inches(1.4), Inches(6.2), Inches(5.2),
    title="▶  고급 물리 거동 및 관성 튜닝", title_size=16, bg=C_CARD_BG, bullet="▸")

section_bar(sl8, "WHTools Drop Simulator  ·  모델 파라미터 및 고급 튜닝")


# ==============================================================================
# Slide 9 — 사용 설명서 3: OpenRadioss FEM 해석 및 후처리
# ==============================================================================
sl9 = blank_slide(prs)
fill_bg(sl9, C_BG)
add_header(sl9, "사용 방법  —  OpenRadioss FEM 해석 및 후처리", "08. OPERATION MANUAL  —  FEM LINKAGE")

add_card_box(sl9,
    ["최대 충격 스냅샷 캡처:\n- MuJoCo 시뮬레이션 진행 중, 가장 변형이 크거나 상세한 강도 평가가 요구되는 충격 프레임(또는 특정 관심 시간대)의 위치/자세를 캡처",
     "OpenRadioss 모델 자동 빌드:\n- [Export to Radioss] 기능 실행 시 캡처된 자세 및 조건을 반영한 OpenRadioss 전용 해석 입력 파일(.rad, .inc)을 디렉토리에 자동 생성",
     "Ground Penetration 사전 보정:\n- FEM 해석 수렴 에러를 방지하기 위해 지면과의 초기 침투(Penetration) 여부를 사전 검출(0.5 mm 이하)하고 z축 방향 오프셋 자동 보정"],
    Inches(0.4), Inches(1.4), Inches(6.0), Inches(5.2),
    title="▶  MuJoCo 상태 기반 FEM 모델 자동 빌드", title_size=16, bg=C_CARD_BG, bullet="•")

add_card_box(sl9,
    ["VTKHDF 시계열 출력:\n- JAX Kirchhoff 해석기에서 계산된 응력과 변위의 시계열 데이터를 VTKHDF 포맷으로 출력하여 ParaView에서 고속 로드 및 분석",
     "GLB 3D 내보내기:\n- 최대 응력이 작용하는 형상의 Von-Mises 응력 분포를 GLB 포맷으로 변환해 웹 브라우저나 PT용 3D 뷰어에서 회전 및 확대 분석 가능",
     "배치 Excel 리포트 자동 생성:\n- 각 시나리오별 Cushion, Box, Chassis의 최대 응력값, 작용 시간, 안전율(Safety Factor) 요약 데이터 테이블을 원클릭 리포팅"],
    Inches(6.7), Inches(1.4), Inches(6.2), Inches(5.2),
    title="▶  다양한 데이터 익스포트 및 분석 리포팅", title_size=16, bg=C_CARD_BG, bullet="•")

section_bar(sl9, "WHTools Drop Simulator  ·  고정밀 FEM 연동 및 3D 후처리")


# ==============================================================================
# Slide 10 — 주요 성과 및 차별성
# ==============================================================================
sl10 = blank_slide(prs)
fill_bg(sl10, C_BG)
add_header(sl10, "주요 성과 및 차별성", "09. VALUE & PERFORMANCE COMPARISON")

# 비교 테이블 구현 (add_rect와 add_text를 활용한 그리드 형태 구성)
# 테이블 좌표: Left 0.4, Top 1.4, Width 12.5, Height 4.0
headers_tab = ["평가 항목", "WHTools 통합 솔버 (MuJoCo/JAX)", "기존 상용 FEM 솔버 (Abaqus / LS-DYNA)"]
col_widths  = [Inches(2.5), Inches(5.0), Inches(5.0)]
row_h       = Inches(0.7)
start_x     = Inches(0.4)
start_y     = Inches(1.4)

# 헤더 그리기
cur_x = start_x
for idx, text in enumerate(headers_tab):
    add_rect(sl10, cur_x, start_y, col_widths[idx], row_h, C_PRIMARY)
    add_text(sl10, text, cur_x + Inches(0.15), start_y + Inches(0.15), col_widths[idx] - Inches(0.3), row_h - Inches(0.3),
             size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    cur_x += col_widths[idx]

# 로우 데이터
row_data = [
    ("평균 해석 시간", "시나리오별 약 1~2분 소요 (초고속)\n- 22가지 배치 전체 수행 시 약 30분 이내", "단일 시나리오당 수 시간 ~ 수일 소요\n- 대규모 격자 생성 및 접촉 계산 부하"),
    ("전처리 편의성", "마커 궤적 기반 무치수 자동 정렬 지원\n- CAD 원본 유실 시에도 마커 데이터만으로 해석 가능", "복잡한 CAD 기하학 정리(Cleaning) 필수\n- 격자(Mesh) 생성, 접촉 쌍(Contact Pair) 수동 수립"),
    ("배치 시나리오", "ISTA 22가지 시험 규격 원클릭 병렬 자동 빌드\n- 다중 작업자 환경에서 CPU 헤드리스 실행", "각 케이스별 경계 조건 수동 지정 필요\n- 시나리오가 늘어날수록 전처리 공수 급증"),
    ("도입 및 유지 비용", "오픈소스 엔진(MuJoCo, OpenRadioss) 활용\n- 라이선스 비용 Zero, 클라우드 스케일링 용이", "솔버 코어(Core) 라이센스 구매 비용 과다\n- 동시 해석 개수 제한으로 병목 현상 발생")
]

for row_idx, row_vals in enumerate(row_data):
    y = start_y + row_h + row_idx * Inches(1.0)
    bg_row = C_ACCENT if row_idx % 2 == 0 else C_WHITE
    cur_x = start_x
    for col_idx, text in enumerate(row_vals):
        add_rect(sl10, cur_x, y, col_widths[col_idx], Inches(1.0), bg_row, border_color=C_BORDER)
        # 1열은 볼드체 처리
        is_bold = (col_idx == 0)
        text_color = C_PRIMARY if col_idx == 1 else C_TEXT
        # 정렬
        align_type = PP_ALIGN.CENTER if col_idx == 0 else PP_ALIGN.LEFT
        add_text(sl10, text, cur_x + Inches(0.15), y + Inches(0.12), col_widths[col_idx] - Inches(0.3), Inches(0.8),
                 size=11, bold=is_bold, color=text_color, align=align_type)
        cur_x += col_widths[col_idx]

# 하단 요약 문구
add_rect(sl10, Inches(0.4), Inches(5.9), W - Inches(0.8), Inches(0.8), C_ACCENT, border_color=C_BORDER)
add_text(sl10, "📌 WHTools는 제품 설계 초기 단계의 신속한 스크리닝 및 최적화를 주도하며, 상세 변형 및 최종 신뢰성 검증 단계에서만 선택적으로 OpenRadioss FEM과 연계하여 해석 프로세스의 효율을 극대화합니다.",
         Inches(0.6), Inches(6.0), W - Inches(1.2), Inches(0.6), size=11, color=C_TEXT, bold=True)

section_bar(sl10, "WHTools Drop Simulator  ·  가치 제안 및 비교 분석")


# ==============================================================================
# Slide 11 — 마무리 및 Q&A
# ==============================================================================
sl11 = blank_slide(prs)
fill_bg(sl11, C_PRIMARY)

add_rect(sl11, 0, 0, Inches(0.15), H, C_SECONDARY)

# 로고 삽입
if LOGO.exists():
    sl11.shapes.add_picture(str(LOGO), Inches(4.9), Inches(0.6), Inches(3.5), Inches(3.5))

add_text(sl11, "WHTools TV Drop Motion Simulator v6.0", Inches(0.5), Inches(1.8), Inches(12.33), Inches(0.8),
         size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl11, "포장 설계 혁신을 위한 낙하 시뮬레이션 자동화 플랫폼",
         Inches(0.5), Inches(2.7), Inches(12.33), Inches(0.6), size=20, color=C_ACCENT, align=PP_ALIGN.CENTER)

add_rect(sl11, Inches(2.0), Inches(3.6), Inches(9.33), Inches(0.04), C_SECONDARY)

add_text(sl11, "기술 피드백 및 개발 협력 문의: whbest.lee@gmail.com",
         Inches(0.5), Inches(3.9), Inches(12.33), Inches(0.5), size=15, color=C_ACCENT, align=PP_ALIGN.CENTER)

# 기술 스택 상세 나열
stack_txt = (
    "■ Physics Core: MuJoCo 3.x · SciPy  |  ■ Structural Solver: JAX (Kirchhoff plate theory) · NumPy\n"
    "■ Graphic / UI: PySide6 (Qt6) · PyQtGraph · PyVista (3D Rendering)\n"
    "■ FEA Integration: OpenRadioss (Starter/Engine) · ParaView · LS-PrePost\n"
    "■ Data format: HDF5 (VTKHDF) · Pickle · GLB · CSV"
)
add_text(sl11, stack_txt, Inches(1.0), Inches(4.7), Inches(11.33), Inches(1.5),
         size=12, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(sl11, "Build #115  ·  MuJoCo · OpenRadioss · PySide6  ·  WHTools Engineering Team",
         Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.45),
         size=11, color=C_GRAY, align=PP_ALIGN.CENTER)

# ── 저장 ──────────────────────────────────────────────────
out = Path(r"C:\Users\GOODMAN\WHToolsBox\WHTools_DropSimulator_Manual.pptx")
prs.save(str(out))
print(f"[SUCCESS] Saved PPTX manual to -> {out}")
